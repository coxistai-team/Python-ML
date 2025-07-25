from flask import Flask, request, jsonify, send_file, render_template_string
import os
import uuid
import traceback
from werkzeug.utils import secure_filename
from datetime import datetime
from flask_cors import CORS
import json
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64

# Import the PowerPoint generator functions
from modules.pptfinal import generate_ai_content, create_powerpoint

app = Flask(__name__)
CORS(app, resources={
    r"/*": {
        "origins": os.getenv("ALLOWED_ORIGINS", "*").split(","),
        "methods": ["GET", "POST", "PUT", "DELETE", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization"],
        "supports_credentials": True
    }
})

# Configuration
UPLOAD_FOLDER = 'generated_ppts'
PRESENTATIONS_FOLDER = 'presentations'
JSON_FOLDER = 'presentation_json'
UPLOADS_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pptx'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PRESENTATIONS_FOLDER, exist_ok=True)
os.makedirs(JSON_FOLDER, exist_ok=True)
os.makedirs(UPLOADS_FOLDER, exist_ok=True)

# Your API key for the AI content generation
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

# In-memory storage for presentations (in production, use a database)
presentations_db = {}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def generate_unique_filename(topic):
    """Generate a unique filename based on topic and timestamp"""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_topic = secure_filename(topic)[:50]
    unique_id = str(uuid.uuid4())[:8]
    return f"{safe_topic}_{timestamp}_{unique_id}.pptx"

def extract_image_from_shape(shape):
    """Extract image data from a shape and return base64 encoded string"""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = image.blob
            # Convert to base64
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            # Determine image format
            image_format = 'png'  # Default
            if image_bytes.startswith(b'\xff\xd8'):
                image_format = 'jpeg'
            elif image_bytes.startswith(b'\x89PNG'):
                image_format = 'png'
            elif image_bytes.startswith(b'GIF'):
                image_format = 'gif'
                
            return f"data:image/{image_format};base64,{image_base64}"
    except Exception as e:
        print(f"Error extracting image: {e}")
    return None

def get_text_color_from_run(run):
    """Extract text color from a text run"""
    try:
        if hasattr(run.font, 'color') and run.font.color.type is not None:
            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                rgb = run.font.color.rgb
                return f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
    except:
        pass
    return "#000000"

def get_slide_background_color(slide):
    """Extract slide background color"""
    try:
        if hasattr(slide, 'background') and hasattr(slide.background, 'fill'):
            fill = slide.background.fill
            if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                rgb = fill.fore_color.rgb
                return f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
    except:
        pass
    return "#ffffff"

def extract_ppt_to_json(ppt_path):
    """Extract PowerPoint content to JSON format for accurate web rendering"""
    try:
        prs = Presentation(ppt_path)
        presentation_json = {
            "metadata": {
                "title": "",
                "slide_count": len(prs.slides),
                "created_at": datetime.now().isoformat(),
                "theme": "gamma_style"
            },
            "slides": []
        }
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_bg_color = get_slide_background_color(slide)
            
            slide_data = {
                "id": f"slide_{slide_idx}",
                "slide_number": slide_idx + 1,
                "layout_type": "title" if slide_idx == 0 else "content",
                "background": {
                    "type": "color",
                    "color": slide_bg_color
                },
                "elements": []
            }
            
            # Extract text elements and images
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Determine if this is a bullet list
                    is_bullet_list = False
                    bullet_items = []
                    
                    if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                        # Check if any paragraph has bullet formatting
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                # Simple bullet detection
                                text = para.text.strip()
                                if text.startswith('•') or text.startswith('-') or text.startswith('*'):
                                    is_bullet_list = True
                                    bullet_items.append(text.lstrip('•-* '))
                                elif is_bullet_list:
                                    bullet_items.append(text)
                    
                    # Create element based on type
                    if is_bullet_list and bullet_items:
                        element = {
                            "type": "bullet_list",
                            "items": bullet_items,
                            "position": {
                                "left": float(shape.left.inches) if hasattr(shape.left, 'inches') else 0,
                                "top": float(shape.top.inches) if hasattr(shape.top, 'inches') else 0,
                                "width": float(shape.width.inches) if hasattr(shape.width, 'inches') else 0,
                                "height": float(shape.height.inches) if hasattr(shape.height, 'inches') else 0
                            },
                            "style": {
                                "font_size": 16,
                                "color": "#333333",
                                "alignment": "left"
                            }
                        }
                    else:
                        # Regular text element
                        element_type = "title" if (slide_idx == 0 and shape.top.inches < 3) or (slide_idx > 0 and shape.top.inches < 2) else "text"
                        
                        element = {
                            "type": element_type,
                            "content": shape.text.strip(),
                            "position": {
                                "left": float(shape.left.inches) if hasattr(shape.left, 'inches') else 0,
                                "top": float(shape.top.inches) if hasattr(shape.top, 'inches') else 0,
                                "width": float(shape.width.inches) if hasattr(shape.width, 'inches') else 0,
                                "height": float(shape.height.inches) if hasattr(shape.height, 'inches') else 0
                            },
                            "style": {
                                "font_size": 32 if element_type == "title" else 16,
                                "font_weight": "bold" if element_type == "title" else "normal",
                                "color": "#333333",
                                "alignment": "center" if element_type == "title" else "left"
                            }
                        }
                    
                    # Try to get actual font properties
                    try:
                        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                if hasattr(run.font, 'size') and run.font.size:
                                    element["style"]["font_size"] = int(run.font.size.pt)
                                if hasattr(run.font, 'bold') and run.font.bold:
                                    element["style"]["font_weight"] = "bold"
                                
                                # Get text color
                                color = get_text_color_from_run(run)
                                element["style"]["color"] = color
                    except:
                        pass
                    
                    slide_data["elements"].append(element)
                
                # Extract images
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data = extract_image_from_shape(shape)
                        if image_data:
                            image_element = {
                                "type": "image",
                                "position": {
                                    "left": float(shape.left.inches),
                                    "top": float(shape.top.inches),
                                    "width": float(shape.width.inches),
                                    "height": float(shape.height.inches)
                                },
                                "src": image_data,
                                "alt": "Slide image"
                            }
                            slide_data["elements"].append(image_element)
                    except Exception as e:
                        print(f"Error processing image: {e}")
            
            presentation_json["slides"].append(slide_data)
        
        return presentation_json
        
    except Exception as e:
        print(f"Error extracting PPT to JSON: {str(e)}")
        traceback.print_exc()
        return None

@app.route('/')
def test():
    return jsonify({
        "status": "Flask is running!",
        "message": "AI Presentation Generator API - Enhanced with JSON Export and PPTX Upload",
        "endpoints": {
            "create": "/create_presentation",
            "upload": "/upload_pptx",
            "update": "/update_slide",
            "download": "/download_ppt/<presentation_id>",
            "get": "/get_presentation/<presentation_id>",
            "json": "/get_presentation_json/<presentation_id>",
            "export": "/export_ppt"
        }
    })

@app.route('/create_presentation', methods=['POST', 'OPTIONS'])
def create_presentation():
    """Create a new presentation with JSON export capability"""
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'OK'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
        response.headers.add('Access-Control-Allow-Methods', 'POST, OPTIONS')
        return response
    
    try:
        print("=== CREATE PRESENTATION REQUEST ===")
        data = request.get_json()
        print(f"Request data: {data}")
        
        if not data:
            return jsonify({
                'success': False,
                'error': 'No JSON data received'
            }), 400
        
        topic = data.get('topic', '').strip()
        num_slides = data.get('slides', 5)
        
        if not topic:
            return jsonify({
                'success': False,
                'error': 'Topic is required'
            }), 400
        
        try:
            num_slides = int(num_slides)
            if num_slides < 1 or num_slides > 20:
                return jsonify({
                    'success': False,
                    'error': 'Number of slides must be between 1 and 20'
                }), 400
        except (ValueError, TypeError):
            return jsonify({
                'success': False,
                'error': 'Invalid number of slides'
            }), 400

        print(f"Creating presentation about '{topic}' with {num_slides} slides")
        
        # Generate AI content
        content = generate_ai_content(topic, num_slides, OPENROUTER_API_KEY)
        
        if not content:
            return jsonify({
                'success': False,
                'error': 'Failed to generate presentation content'
            }), 500
        
        # Create PowerPoint file
        ppt_path = create_powerpoint(content, topic)
        
        if not ppt_path or not os.path.exists(ppt_path):
            return jsonify({
                'success': False,
                'error': 'Failed to create PowerPoint file'
            }), 500
        
        # Extract to JSON for web rendering
        presentation_json = extract_ppt_to_json(ppt_path)
        
        # Generate presentation ID
        presentation_id = f"pres_{int(datetime.now().timestamp())}_{str(uuid.uuid4())[:8]}"
        
        # Store presentation in memory
        presentation_data = {
            'id': presentation_id,
            'topic': topic,
            'created_at': datetime.now().isoformat(),
            'updated_at': datetime.now().isoformat(),
            'slide_count': len(content),
            'ppt_path': ppt_path,
            'json_data': presentation_json
        }
        
        presentations_db[presentation_id] = presentation_data
        
        # Save to file as backup
        presentation_file = os.path.join(PRESENTATIONS_FOLDER, f"{presentation_id}.json")
        with open(presentation_file, 'w') as f:
            json.dump(presentation_data, f, indent=2)
        
        # Save JSON representation separately
        if presentation_json:
            json_file = os.path.join(JSON_FOLDER, f"{presentation_id}_structure.json")
            with open(json_file, 'w') as f:
                json.dump(presentation_json, f, indent=2)
        
        response_data = {
            'success': True,
            'presentation_id': presentation_id,
            'topic': topic,
            'slide_count': len(content),
            'message': f'Successfully created {len(content)} slides about {topic}',
            'created_at': presentation_data['created_at'],
            'json_available': presentation_json is not None
        }
        
        print("=== PRESENTATION CREATED SUCCESSFULLY ===")
        return jsonify(response_data)
        
    except Exception as e:
        print(f"Error in create_presentation: {str(e)}")
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': f'Internal server error: {str(e)}'
        }), 500

@app.route('/get_presentation_json/<presentation_id>', methods=['GET'])
def get_presentation_json(presentation_id):
    """Get presentation JSON structure for web rendering"""
    try:
        print(f"Getting JSON for presentation: {presentation_id}")
        
        # Try in-memory first
        if presentation_id in presentations_db:
            presentation = presentations_db[presentation_id]
            if 'json_data' in presentation and presentation['json_data']:
                return jsonify({
                    'success': True,
                    'presentation_id': presentation_id,
                    'json_data': presentation['json_data']
                })
        
        # Fallback to file
        json_file = os.path.join(JSON_FOLDER, f"{presentation_id}_structure.json")
        if os.path.exists(json_file):
            with open(json_file, 'r') as f:
                json_data = json.load(f)
            
            return jsonify({
                'success': True,
                'presentation_id': presentation_id,
                'json_data': json_data
            })
        
        # Try to load presentation and extract JSON
        presentation_file = os.path.join(PRESENTATIONS_FOLDER, f"{presentation_id}.json")
        if os.path.exists(presentation_file):
            with open(presentation_file, 'r') as f:
                presentation_data = json.load(f)
            
            # If we have the PPT path, extract JSON
            if 'ppt_path' in presentation_data and os.path.exists(presentation_data['ppt_path']):
                json_data = extract_ppt_to_json(presentation_data['ppt_path'])
                if json_data:
                    return jsonify({
                        'success': True,
                        'presentation_id': presentation_id,
                        'json_data': json_data
                    })
        
        return jsonify({'success': False, 'error': 'Presentation JSON not found'}), 404
        
    except Exception as e:
        print(f"Error getting presentation JSON: {str(e)}")
        traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/update_slide', methods=['PUT', 'OPTIONS'])
def update_slide():
    """Update a specific slide in a presentation"""
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'OK'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
        response.headers.add('Access-Control-Allow-Methods', 'PUT, OPTIONS')
        return response
    
    try:
        data = request.get_json()
        presentation_id = data.get('presentation_id')
        slide_id = data.get('slide_id')
        slide_data = data.get('slide_data')
        
        if not all([presentation_id, slide_id, slide_data]):
            return jsonify({
                'success': False,
                'error': 'Missing required fields: presentation_id, slide_id, slide_data'
            }), 400
        
        # Load presentation
        if presentation_id not in presentations_db:
            presentation_file = os.path.join(PRESENTATIONS_FOLDER, f"{presentation_id}.json")
            if os.path.exists(presentation_file):
                with open(presentation_file, 'r') as f:
                    presentations_db[presentation_id] = json.load(f)
            else:
                return jsonify({
                    'success': False,
                    'error': 'Presentation not found'
                }), 404
        
        presentation = presentations_db[presentation_id]
        
        # Update the slide in JSON data
        if 'json_data' in presentation and presentation['json_data']:
            slides = presentation['json_data']['slides']
            for i, slide in enumerate(slides):
                if slide['id'] == slide_id:
                    # Update slide elements
                    slides[i] = slide_data
                    break
        
        # Update timestamp
        presentation['updated_at'] = datetime.now().isoformat()
        
        # Save to file
        presentation_file = os.path.join(PRESENTATIONS_FOLDER, f"{presentation_id}.json")
        with open(presentation_file, 'w') as f:
            json.dump(presentation, f, indent=2)
        
        return jsonify({
            'success': True,
            'message': 'Slide updated successfully'
        })
        
    except Exception as e:
        print(f"Error updating slide: {str(e)}")
        return jsonify({
            'success': False,
            'error': f'Error updating slide: {str(e)}'
        }), 500

@app.route('/export_ppt', methods=['POST', 'OPTIONS'])
def export_ppt():
    """Export presentation as PowerPoint or PDF"""
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'OK'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
        response.headers.add('Access-Control-Allow-Methods', 'POST, OPTIONS')
        return response
    
    try:
        data = request.get_json()
        presentation_id = data.get('presentationId')
        export_format = data.get('format', 'pptx')
        
        if not presentation_id:
            return jsonify({
                'success': False,
                'error': 'presentation_id is required'
            }), 400
        
        # Load presentation
        if presentation_id not in presentations_db:
            presentation_file = os.path.join(PRESENTATIONS_FOLDER, f"{presentation_id}.json")
            if os.path.exists(presentation_file):
                with open(presentation_file, 'r') as f:
                    presentations_db[presentation_id] = json.load(f)
            else:
                return jsonify({
                    'success': False,
                    'error': 'Presentation not found'
                }), 404
        
        presentation = presentations_db[presentation_id]
        
        if export_format == 'pptx':
            # Return existing PowerPoint file
            if 'ppt_path' in presentation and os.path.exists(presentation['ppt_path']):
                return send_file(
                    presentation['ppt_path'],
                    as_attachment=True,
                    download_name=f"{presentation['topic']}.pptx",
                    mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
            else:
                return jsonify({
                    'success': False,
                    'error': 'PowerPoint file not found'
                }), 404
        
        elif export_format == 'pdf':
            # Create PDF export
            return export_as_pdf(presentation)
        
        else:
            return jsonify({
                'success': False,
                'error': 'Unsupported export format'
            }), 400
            
    except Exception as e:
        print(f"Export error: {str(e)}")
        return jsonify({
            'success': False,
            'error': f'Export error: {str(e)}'
        }), 500

def export_as_pdf(presentation):
    """Export presentation as PDF"""
    try:
        topic = presentation['topic']
        
        # Create PDF in memory
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'PPTTitle',
            parent=styles['Heading1'],
            fontSize=28,
            spaceAfter=20,
            textColor=HexColor('#1f2937'),
            alignment=1,
            fontName='Helvetica-Bold'
        )
        
        slide_title_style = ParagraphStyle(
            'PPTSlideTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=15,
            textColor=HexColor('#2563eb'),
            alignment=1,
            fontName='Helvetica-Bold'
        )
        
        content_style = ParagraphStyle(
            'PPTContent',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=12,
            leftIndent=20,
            fontName='Helvetica',
            leading=16
        )
        
        story = []
        
        # Title page
        story.append(Spacer(1, 1*inch))
        story.append(Paragraph(topic, title_style))
        story.append(Paragraph("AI Generated Presentation", styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Created on {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
        story.append(PageBreak())
        
        # Content slides from JSON data
        if 'json_data' in presentation and presentation['json_data']:
            slides = presentation['json_data']['slides']
            for i, slide in enumerate(slides):
                story.append(Paragraph(f"Slide {i+1}", styles['Normal']))
                
                # Find title element
                title_text = f"Slide {i+1}"
                for element in slide.get('elements', []):
                    if element.get('type') == 'title':
                        title_text = element.get('content', title_text)
                        break
                
                story.append(Paragraph(title_text, slide_title_style))
                
                # Add content elements
                for element in slide.get('elements', []):
                    if element.get('type') == 'text':
                        story.append(Paragraph(element.get('content', ''), content_style))
                    elif element.get('type') == 'bullet_list':
                        for item in element.get('items', []):
                            story.append(Paragraph(f"• {item}", content_style))
                
                story.append(Spacer(1, 0.3*inch))
                if i < len(slides) - 1:
                    story.append(PageBreak())
        
        doc.build(story)
        buffer.seek(0)
        
        return send_file(
            buffer,
            as_attachment=True,
            download_name=f"{topic}.pdf",
            mimetype='application/pdf'
        )
        
    except Exception as e:
        print(f"PDF export error: {str(e)}")
        return jsonify({
            'success': False,
            'error': f'PDF export error: {str(e)}'
        }), 500

if __name__ == '__main__':
    print("Starting Enhanced AI Presentation Generator API...")
    print("API will be available at: http://localhost:5002")
    print("\nKey Features:")
    print("- AI-powered presentation generation")
    print("- JSON extraction for accurate web rendering")
    print("- PowerPoint and PDF export")
    print("- Slide editing and updates")
    print("\nMain Endpoints:")
    print("POST /create_presentation - Create new presentation")
    print("GET /get_presentation_json/<id> - Get JSON for web rendering")
    print("PUT /update_slide - Update slide content")
    print("POST /export_ppt - Export as PowerPoint or PDF")
    
    app.run(host='0.0.0.0', port=5002, debug=True)
